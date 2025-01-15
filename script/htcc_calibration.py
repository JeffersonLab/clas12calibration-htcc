#!/usr/bin/env python3

import numpy as np
import glob
import os
import sys
import pandas as pd
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from datetime import date as datetime_date
import argparse

def deleteLeadingZeros(string):
    """Remove leading zeros from a string."""
    return string.lstrip('0')

def apply_color(val):
    """Determine color based on value threshold."""
    if abs(val) > 10:
        return 'red'
    elif 5 < abs(val) <= 10:
        return 'yellow'
    else:
        return 'white'

def crop_image(image_path):
    """Crop image to remove white space."""
    with Image.open(image_path) as img:
        img = img.convert("RGB")
        bbox = img.getbbox()
        if bbox:
            img_cropped = img.crop(bbox)
            img_cropped.save(image_path)

def stackImagesVertically(listImages, savePath):
    """Stack multiple images vertically and save."""
    images = [Image.open(x) for x in listImages]
    widths, heights = zip(*(i.size for i in images))

    total_height = sum(heights)
    max_width = max(widths)

    new_im = Image.new('RGB', (max_width, total_height))

    y_offset = 0
    for im in images:
        new_im.paste(im, (0, y_offset))
        y_offset += im.size[1]

    new_im.save(savePath)

def stackImagesHorizontally(listImages, savePath):
    """Stack multiple images horizontally and save."""
    images = [Image.open(x) for x in listImages]
    widths, heights = zip(*(i.size for i in images))

    total_width = sum(widths)
    max_height = max(heights)

    new_im = Image.new('RGB', (total_width, max_height))

    x_offset = 0
    for im in images:
        new_im.paste(im, (x_offset, 0))
        x_offset += im.size[0]

    new_im.save(savePath)

def process_directories(superdirectories, dates):
    """Process all directories and create visualization images."""
    for superdir in superdirectories:
        pathdirs = glob.glob(superdir + "*/")
        runNums = []

        for i in pathdirs:
            tmp = os.path.basename(os.path.dirname(i))
            num = int(deleteLeadingZeros(tmp))
            runNums.append(num)

        count = 0
        for i in pathdirs:
            count += 1
            per = 100 * (count / len(pathdirs))
            path = os.path.dirname(i)
            tmp = os.path.basename(path)
            num = int(deleteLeadingZeros(tmp))

            for subDate in dates:
                savePath = f"{path}/{subDate}/"
                if not os.path.exists(savePath):
                    print(f"Skipping {savePath} as it doesn't exist.")
                    continue

                print(f"Processing {savePath} for run number {num}")

                # Process gain data
                df_gain_path = f"{savePath}correctionFactor_NphePMT{num}.dat"
                if os.path.exists(df_gain_path):
                    df_gain = pd.read_csv(df_gain_path)
                    df_gain = df_gain.drop(['StdDev', 'AverageNPE', "gainTableValue", "NewNphe", "Difference"], axis=1)

                    # Create table visualization for gain data
                    cell_colors = [[apply_color(row["PercentChange"]) if col == "PercentChange" else 'white' 
                                  for col in df_gain.columns] for _, row in df_gain.iterrows()]

                    fig, ax = plt.subplots(figsize=(10, 15))
                    ax.axis('tight')
                    ax.axis('off')
                    table = ax.table(cellText=df_gain.values, colLabels=df_gain.columns, 
                                   cellColours=cell_colors, loc='center', cellLoc='center')

                    table.auto_set_font_size(False)
                    table.set_fontsize(8)
                    table.auto_set_column_width(col=list(range(len(df_gain.columns))))

                    save_path = f"{savePath}CompareCSV_Run{num}_Gain.png"
                    plt.savefig(save_path, bbox_inches='tight', pad_inches=0)
                    plt.close()
                    crop_image(save_path)
                else:
                    print(f"Gain file not found: {df_gain_path}")

                # Process time data
                df_time_path = f"{savePath}correctionFactor_TimePMT{num}.dat"
                if os.path.exists(df_time_path):
                    df_time = pd.read_csv(df_time_path)
                    df_time = df_time.drop(['CurrentTime', "Difference"], axis=1)

                    cell_colors = [[apply_color(row["PercentChange"]) if col == "PercentChange" else 'white' 
                                  for col in df_time.columns] for _, row in df_time.iterrows()]

                    fig, ax = plt.subplots(figsize=(10, 15))
                    ax.axis('tight')
                    ax.axis('off')
                    table = ax.table(cellText=df_time.values, colLabels=df_time.columns,
                                   cellColours=cell_colors, loc='center', cellLoc='center')

                    table.auto_set_font_size(False)
                    table.set_fontsize(8)
                    table.auto_set_column_width(col=list(range(len(df_time.columns))))

                    save_path = f"{savePath}CompareCSV_Run{num}_Time.png"
                    plt.savefig(save_path, bbox_inches='tight', pad_inches=0)
                    plt.close()
                    crop_image(save_path)
                else:
                    print(f"Time file not found: {df_time_path}")

                # Create combined image
                if os.path.exists(f"{savePath}CompareCSV_Run{num}_Gain.png") and \
                   os.path.exists(f"{savePath}CompareCSV_Run{num}_Time.png"):
                    stackImagesHorizontally([f"{savePath}CompareCSV_Run{num}_Gain.png",
                                           f"{savePath}CompareCSV_Run{num}_Time.png"],
                                          f"{savePath}CompareCSV_Run{num}_Combo.png")

            if (per % 5 == 0):
                print(f"PERCENT {int(per)}% COMPLETE")

def addImageSlide(prs, slide, image_url):
    """Add an image to a PowerPoint slide with proper scaling."""
    im = Image.open(image_url)
    image_width, image_height = im.size

    slide_width_pixels = int(prs.slide_width / 914400 * 96)
    slide_height_pixels = int(prs.slide_height / 914400 * 96)

    width_scale = slide_width_pixels / image_width
    height_scale = slide_height_pixels / image_height
    scale = min(width_scale, height_scale)

    new_width = int(image_width * scale)
    new_height = int(image_height * scale)
    im = im.resize((new_width, new_height), Image.Resampling.LANCZOS)

    temp_image_path = image_url.replace(".png", "_resized.png")
    im.save(temp_image_path)

    left = (prs.slide_width - Inches(new_width / 96)) / 2
    top = (prs.slide_height - Inches(new_height / 96)) / 2
    slide.shapes.add_picture(temp_image_path, left, top, 
                           width=Inches(new_width / 96), 
                           height=Inches(new_height / 96))

    os.remove(temp_image_path)  # Clean up temporary file
    return slide

def create_slides(superdirectories, dates, output_dir, slide_title, author_name):
    """Create PowerPoint presentation with calibration results."""
    all_paths = []
    for superdir in superdirectories:
        for date in dates:
            pathdirs = glob.glob(os.path.join(superdir, f"*/{date}/"))
            for path in pathdirs:
                run_num = int(deleteLeadingZeros(os.path.basename(os.path.dirname(os.path.dirname(path)))))
                all_paths.append((run_num, date, path))

    all_paths.sort(key=lambda x: (x[0], x[1]))

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = slide_title
    subtitle.text = f"{author_name}\n{datetime_date.today():%m-%d-%Y}"

    # Create slides for each run
    for run_num, date_str, savePath in all_paths:
        # Run title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = f"PLOTS FOR \nRUN NUMBER {run_num}\n{date_str}"

        # Create combined images
        nphe_image = savePath + f"nphePMT_ZOOM_{run_num}.png"
        time_image = savePath + f"timePMT{run_num}.png"
        image_allChannels = savePath + f"Combo_GainTime_{run_num}.png"
        
        if os.path.exists(nphe_image) and os.path.exists(time_image):
            stackImagesVertically([nphe_image, time_image], image_allChannels)

            slide = prs.slides.add_slide(prs.slide_layouts[8])
            title = slide.shapes.title
            title.text = f"Run {run_num} - {date_str} All Channels"
            slide = addImageSlide(prs, slide, image_allChannels)

        image_allCompare = savePath + f"CompareCSV_Run{run_num}_Combo.png"
        if os.path.exists(image_allCompare):
            slide = prs.slides.add_slide(prs.slide_layouts[8])
            title = slide.shapes.title
            title.text = f"Run {run_num} - {date_str} Comparison"
            slide = addImageSlide(prs, slide, image_allCompare)

    # Create a descriptive filename using the dates
    date_str = "_".join(dates).replace("-", "")  # Convert "08-Jan-2025" to "08Jan2025"
    filename = f"HTCC_Calibrations_{date_str}.pptx"
    full_path = os.path.join(output_dir, filename)
    
    # Save the presentation
    prs.save(full_path)
    print(f"Presentation saved as: {full_path}")

def main():
    parser = argparse.ArgumentParser(description='Process HTCC calibration data and create presentation.')
    parser.add_argument('--directories', nargs='+', required=True,
                        help='List of super directories containing the run data')
    parser.add_argument('--dates', nargs='+', required=True,
                        help='List of dates to process (format: DD-Mon-YYYY)')
    parser.add_argument('--output', required=True,
                        help='Output directory for the presentation')
    parser.add_argument('--title', required=True,
                        help='Title for the presentation')
    parser.add_argument('--author', required=True,
                        help='Author name for the presentation')
    parser.add_argument('--process-only', action='store_true',
                        help='Only process directories without creating slides')
    parser.add_argument('--slides-only', action='store_true',
                        help='Only create slides without processing directories')

    args = parser.parse_args()

    if not args.slides_only:
        process_directories(args.directories, args.dates)
    
    if not args.process_only:
        create_slides(args.directories, args.dates, args.output, args.title, args.author)

if __name__ == "__main__":
    main()